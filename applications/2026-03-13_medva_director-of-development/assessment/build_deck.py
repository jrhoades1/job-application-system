"""Build MEDVA Pulse Portal proposal slide deck — clean dark theme."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Colors (dark theme) ---
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
DARK_BG = RGBColor(0x1E, 0x29, 0x3B)
TEAL = RGBColor(0x00, 0xBC, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEXT = RGBColor(0xE0, 0xE0, 0xE0)
MED_TEXT = RGBColor(0xA0, 0xA8, 0xB8)
DARK_TEXT = RGBColor(0x2C, 0x2C, 0x2C)
ACCENT_ORANGE = RGBColor(0xFF, 0xA7, 0x26)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_RED = RGBColor(0xEF, 0x53, 0x50)
TABLE_HEADER = RGBColor(0x26, 0x32, 0x38)
TABLE_ROW_1 = RGBColor(0x2A, 0x3A, 0x50)
TABLE_ROW_2 = RGBColor(0x23, 0x33, 0x48)
CALLOUT_BG = RGBColor(0x2A, 0x3C, 0x55)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = Inches(13.333)


def bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = DARK_BG


def bar(slide):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SW, Inches(0.06), )
    s.fill.solid()
    s.fill.fore_color.rgb = TEAL
    s.line.fill.background()


def page_num(slide, num, total):
    tx = slide.shapes.add_textbox(Inches(11.8), Inches(7.0), Inches(1.2), Inches(0.4))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num} / {total}"
    p.font.size = Pt(11)
    p.font.color.rgb = MED_TEXT
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.RIGHT


def title_text(slide, text, y=Inches(0.3)):
    tx = slide.shapes.add_textbox(Inches(0.8), y, Inches(10), Inches(0.8))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"
    # underline bar
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), y + Inches(0.65), Inches(2.5), Inches(0.05))
    s.fill.solid()
    s.fill.fore_color.rgb = TEAL
    s.line.fill.background()


def subtitle(slide, text, y=Inches(1.15)):
    tx = slide.shapes.add_textbox(Inches(0.8), y, Inches(11), Inches(0.5))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = TEAL
    p.font.italic = True
    p.font.name = "Calibri"


def txt(slide, text, x, y, w, h, size=16, color=LIGHT_TEXT, bold=False, align=PP_ALIGN.LEFT):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return tx


def bullets(slide, items, x, y, w, h, size=15, color=LIGHT_TEXT, spacing=1.3):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        p.line_spacing = spacing
        if ": " in item:
            pre, rest = item.split(": ", 1)
            r1 = p.add_run()
            r1.text = pre + ": "
            r1.font.size = Pt(size)
            r1.font.color.rgb = TEAL
            r1.font.bold = True
            r1.font.name = "Calibri"
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.color.rgb = color
            r2.font.name = "Calibri"
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.name = "Calibri"


def table(slide, rows, x, y, w, h, col_widths=None, size=14):
    ts = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
    t = ts.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            t.columns[i].width = cw
    for ri, row in enumerate(rows):
        for ci, cell_text in enumerate(row):
            cell = t.cell(ri, ci)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = str(cell_text)
            r.font.size = Pt(size)
            r.font.name = "Calibri"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri == 0:
                r.font.bold = True
                r.font.color.rgb = TEAL
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER
            else:
                r.font.color.rgb = LIGHT_TEXT
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_1 if ri % 2 == 1 else TABLE_ROW_2


def callout(slide, text, x, y, w, color=ACCENT_ORANGE):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, Inches(0.6))
    s.fill.solid()
    s.fill.fore_color.rgb = CALLOUT_BG
    s.line.color.rgb = color
    s.line.width = Pt(2)
    tx = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.1), w - Inches(0.6), Inches(0.45))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(15)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Calibri"


TOTAL = 14

# ============================================================
# SLIDE 1: Title
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)

txt(sl, "MEDVA PULSE PORTAL", Inches(0.8), Inches(1.5), Inches(10), Inches(0.8),
    size=44, color=WHITE, bold=True)

s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.5), Inches(3), Inches(0.05))
s.fill.solid()
s.fill.fore_color.rgb = TEAL
s.line.fill.background()

txt(sl, "Technical & Organizational Proposal", Inches(0.8), Inches(2.8), Inches(10), Inches(0.5),
    size=22, color=TEAL)

txt(sl, "Prepared by: Jimmy Rhoades\nIncoming Director of Development\nApril 8, 2026",
    Inches(0.8), Inches(3.7), Inches(6), Inches(1.2), size=16, color=MED_TEXT)

txt(sl, "HIPAA-First Design", Inches(0.8), Inches(5.8), Inches(4), Inches(0.4),
    size=13, color=MED_TEXT, bold=True)


# ============================================================
# SLIDE 2: Agenda
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
bar(sl)
title_text(sl, "Agenda")
page_num(sl, 2, TOTAL)

agenda = [
    "Technology Stack Recommendation",
    "Team Composition & Hiring Plan",
    "Product Roadmap (MVP \u2192 V1 \u2192 V2)",
    "Architecture Overview",
    "Budget & Vendor Management",
]

for i, item in enumerate(agenda):
    y = Inches(1.7) + Inches(i * 0.95)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(9.5), Inches(0.75))
    s.fill.solid()
    s.fill.fore_color.rgb = CALLOUT_BG
    s.line.fill.background()
    txt(sl, item, Inches(2.0), y + Inches(0.12), Inches(7.5), Inches(0.5),
        size=22, color=WHITE, bold=True)
    txt(sl, str(i + 1), Inches(10.2), y + Inches(0.1), Inches(0.6), Inches(0.5),
        size=26, color=TEAL, bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3: Tech Stack
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
bar(sl)
title_text(sl, "Technology Stack")
subtitle(sl, "Recommended Stack (All HIPAA BAA-eligible)")
page_num(sl, 3, TOTAL)

rows = [
    ["Layer", "Choice", "Key Reason"],
    ["Frontend", "Next.js 15 + TypeScript + Tailwind + shadcn/ui", "Fast dashboards, huge offshore React talent pool"],
    ["Backend", "NestJS (TypeScript) on Fargate", "Full TS parity, built-in RBAC guards, modular monolith"],
    ["Database", "Aurora Serverless v2 (PostgreSQL) + RLS", "Tenant isolation at DB level, auto-scales to zero"],
    ["Auth", "AWS Cognito + NestJS Guards", "MFA, SSO, HIPAA-eligible. Don't reinvent with 5 people."],
    ["AI", "AWS Bedrock (Claude) + OpenSearch Serverless", "Best RAG grounding + instruction-following, HIPAA BAA"],
    ["Infra", "VPC + Fargate + S3 + CloudFront + WAF + KMS", "Private subnets only, Fargate over EKS (no cluster tax)"],
    ["CI/CD", "GitHub Actions + Terraform + Docker", "IaC everything, build/test/deploy on every PR merge"],
]

table(sl, rows, Inches(0.5), Inches(1.6), Inches(12.3), Inches(4.5),
      col_widths=[Inches(1.3), Inches(5.5), Inches(5.5)], size=14)

callout(sl, "Internal tooling: GitHub, Linear, Figma, Notion, Slack, Claude Code (~$370/mo)",
        Inches(0.5), Inches(6.3), Inches(10))


# ============================================================
# SLIDE 4: HIPAA/Compliance
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
bar(sl)
title_text(sl, "HIPAA/SOC 2 Compliance")
subtitle(sl, "Compliance is a foundation, not a feature")
page_num(sl, 4, TOTAL)

items = [
    "BAA Coverage: Every AWS service on the HIPAA BAA list. Non-BAA for PHI = violation.",
    "Encryption: At rest (KMS, AES-256) and in transit (TLS 1.2+). No exceptions.",
    "Network: VPC with private subnets only. Zero public DB access. WAF + GuardDuty.",
    "Access: Zero developer production access. All changes through application layer with RBAC.",
    "Audit: CloudTrail + custom audit tables for all data access. Immutable logs in S3.",
    "Backup: Daily snapshots, 30-day retention. Multi-AZ. S3 cross-region replication.",
    "Experience: Led SOC 2/HITRUST certification at ilumed. Compliance baked in from Day 1 costs a fraction of retrofitting later.",
]

bullets(sl, items, Inches(0.8), Inches(1.6), Inches(11), Inches(4.5), size=16)

callout(sl, "Architecture enforces compliance at the infrastructure level, not just in application code.",
        Inches(0.5), Inches(6.3), Inches(11), color=TEAL)


# ============================================================
# SLIDE 5: Team Composition
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
bar(sl)
title_text(sl, "Team Composition")
subtitle(sl, "4 Hires + Me (Employee #1, writing 60-70% of foundational code)")
page_num(sl, 5, TOTAL)

rows = [
    ["Priority", "Role", "Location", "Salary Range", "Start"],
    ["1", "DevOps & Security Engineer", "US-based", "$155K - $175K", "Week 1"],
    ["2", "Lead Backend Engineer", "Offshore", "$40K - $55K", "Week 1-2"],
    ["3", "Full-Stack Engineer", "Offshore", "$35K - $50K", "Week 2-3"],
    ["4", "UI/UX Engineer", "Offshore", "$25K - $40K", "Month 2-3"],
]

table(sl, rows, Inches(0.5), Inches(1.6), Inches(12.3), Inches(2.8),
      col_widths=[Inches(1.0), Inches(4.0), Inches(2.0), Inches(2.8), Inches(2.5)], size=15)

callout(sl, "Critical first hire: US DevOps & Security (HIPAA foundation must be rock-solid from Day 1)",
        Inches(0.5), Inches(4.7), Inches(11), color=ACCENT_ORANGE)

# Team evolution
txt(sl, "Team Evolution", Inches(0.8), Inches(5.6), Inches(4), Inches(0.4),
    size=18, color=WHITE, bold=True)

evo = [
    "MVP: Me (60-70% backend) + US DevOps + Lead Backend",
    "V1 (AI): I shift to Bedrock/RAG. Add contract AI/ML engineer. Lead Backend steps up.",
    "Later: Offshore team runs autonomously on features within established patterns",
]
bullets(sl, evo, Inches(0.8), Inches(6.0), Inches(11), Inches(1.2), size=14, color=MED_TEXT)


# ============================================================
# SLIDE 6: Interview Process
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
bar(sl)
title_text(sl, "Interview Process by Role")
page_num(sl, 6, TOTAL)

rows = [
    ["Role", "Key Evaluation", "Process"],
    ["DevOps & Security (US)", "HIPAA infra design, IaC",
     "Infra scenario + Terraform take-home + system-design with me + John"],
    ["Lead Backend (Offshore)", "Tenant isolation, RBAC",
     "Live pair-programming on tenant-isolated endpoint + English assessment"],
    ["Full-Stack (Offshore)", "React + API integration",
     "Live coding + code review exercise (find issues in a PR) + portfolio"],
    ["UI/UX (Offshore)", "Figma-to-code, design system",
     "Live dashboard build from Figma + code review + portfolio"],
]

table(sl, rows, Inches(0.5), Inches(1.3), Inches(12.3), Inches(3.0),
      col_widths=[Inches(2.5), Inches(3.0), Inches(6.8)], size=13)

callout(sl, "Offshore experience: 50+ devs across US/Ukraine/Central America (Cognizant). "
        "Built teams from zero at Red Spot and MedQuest. Async-first, 2-3 hr daily overlap, same standards.",
        Inches(0.5), Inches(6.2), Inches(11), color=TEAL)


# ============================================================
# SLIDE 7: Roadmap
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
bar(sl)
title_text(sl, "Product Roadmap")
subtitle(sl, "MVP \u2192 V1 \u2192 V2")
page_num(sl, 7, TOTAL)

# Three boxes
phases = [
    ("MVP", "10 weeks", ACCENT_GREEN,
     "AWS infra + CI/CD via Terraform\nCognito auth + RBAC (5 roles)\n"
     "Portal landing (Client | VA | Corp)\nClient dashboard: KPIs, Cost Efficiency\n"
     "VA portal: profile, schedule, alerts\nCorporate: accounts, revenue, risk alerts\n"
     "Secure messaging (E2E encrypted)\n\u2192 Pilot launch: 10-20% of clients"),
    ("V1", "+14 weeks (~6 mo)", TEAL,
     "Client KB upload + auto-indexing\nMedva Frontline AI (Claude RAG)\n"
     "Chat history + source citations\nAI Extensions (Open Evidence, etc.)\n"
     "Medva Academy (training modules)\nStripe billing (CC + ACH)\nVA pay stubs\n"
     "\u2192 Full client rollout + new revenue"),
    ("V2", "Months 9-14", ACCENT_ORANGE,
     "Matching/assignment engine\nPredictive analytics (client health)\n"
     "EHR integrations (Epic, Athena)\nReact Native mobile for VAs\n"
     "Self-service client onboarding\nAPI platform for enterprise\n"
     "\u2192 Scale & expand"),
]

for i, (name, timeline, color, content) in enumerate(phases):
    x = Inches(0.5) + Inches(i * 4.2)
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(3.9), Inches(4.6))
    s.fill.solid()
    s.fill.fore_color.rgb = CALLOUT_BG
    s.line.color.rgb = color
    s.line.width = Pt(2)

    txt(sl, name, x + Inches(0.2), Inches(1.7), Inches(2), Inches(0.4),
        size=22, color=color, bold=True)
    txt(sl, timeline, x + Inches(0.2), Inches(2.1), Inches(3), Inches(0.3),
        size=13, color=MED_TEXT)

    tx = sl.shapes.add_textbox(x + Inches(0.2), Inches(2.5), Inches(3.5), Inches(3.5))
    tf = tx.text_frame
    tf.word_wrap = True
    for j, line in enumerate(content.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.line_spacing = 1.25
        r = p.add_run()
        r.text = line
        r.font.size = Pt(13)
        r.font.color.rgb = LIGHT_TEXT if not line.startswith("\u2192") else color
        r.font.bold = line.startswith("\u2192")
        r.font.name = "Calibri"

# Bottom row: success, migration, integration, QA
sections = [
    ("Success Criteria", ["10-20% client adoption", "<2s dashboard load", "Zero cross-tenant incidents", "Weekly active pilot usage"]),
    ("Migration", ["Phased migration with feature flags", "Dual-run validation before cutover", "Existing portal operates in parallel"]),
    ("Integration", ["Dedicated module (REST/webhooks)", "Background jobs + retry queues", "External failures isolated from core"]),
    ("QA", ["Automated tests (unit + API + E2E)", "Weekly demo validation"]),
]

for i, (label, items) in enumerate(sections):
    x = Inches(0.5) + Inches(i * 3.2)
    txt(sl, label, x, Inches(6.35), Inches(3), Inches(0.3), size=12, color=TEAL, bold=True)
    for j, item in enumerate(items):
        txt(sl, item, x, Inches(6.6) + Inches(j * 0.2), Inches(3), Inches(0.2),
            size=10, color=MED_TEXT)


# ============================================================
# SLIDE 8: Technical Risks
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
bar(sl)
title_text(sl, "Technical Risks & Mitigations")
page_num(sl, 8, TOTAL)

rows = [
    ["Phase", "Risk", "Mitigation"],
    ["MVP", "Data isolation / HIPAA breach",
     "organization_id on every table + RLS + automated security scans in CI + external pen-test"],
    ["MVP", "Data migration delays",
     "Start analysis Week 1, build scripts in parallel, shadow mode before cutover"],
    ["MVP", "Scope creep",
     "Lock scope with John upfront. Weekly demos. Say no to anything outside pilot."],
    ["V1", "AI cross-tenant leakage",
     "Per-client OpenSearch namespaces + mandatory org_id filter at infra level + nightly automated tests"],
    ["V1", "AI hallucination / quality",
     "RAG grounds responses in docs. Confidence scoring. Low-confidence flagged for human review."],
    ["V1", "AI cost overrun",
     "Per-client usage quotas, prompt caching, cost-efficient model tier first"],
    ["Both", "Offshore ramp-up",
     "US DevOps first as insurance. Async-first: specs, Loom walkthroughs, 2-3 hr daily overlap."],
]

table(sl, rows, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.0),
      col_widths=[Inches(1.0), Inches(3.3), Inches(8.0)], size=14)


# ============================================================
# SLIDE 9: Architecture
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
bar(sl)
title_text(sl, "Architecture Overview")
page_num(sl, 9, TOTAL)

# Visual flow using shapes
layers = [
    ("Clients  |  Virtual Assistants  |  MEDVA Staff", TEAL, Inches(1.2)),
    ("\u25BC", None, Inches(1.85)),
    ("AWS WAF  \u2192  Cognito Auth (MFA + SSO + RBAC)  \u2192  Next.js 15 Frontend", TEAL, Inches(2.1)),
    ("\u25BC", None, Inches(2.75)),
    ("NestJS API (Fargate)  |  RBAC Guards  |  Tenant Context Middleware", NAVY, Inches(3.0)),
    ("\u25BC", None, Inches(3.65)),
]

for text, color, y in layers:
    if color:
        s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(10.3), Inches(0.55))
        s.fill.solid()
        s.fill.fore_color.rgb = CALLOUT_BG
        s.line.color.rgb = color
        s.line.width = Pt(2)
        txt(sl, text, Inches(1.5), y + Inches(0.05), Inches(10.3), Inches(0.45),
            size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    else:
        txt(sl, text, Inches(6.2), y, Inches(1), Inches(0.3),
            size=18, color=MED_TEXT, align=PP_ALIGN.CENTER)

# Data boxes
data_boxes = [
    ("Aurora Serverless v2\nPostgreSQL + RLS", Inches(1.5)),
    ("S3 (KMS Encrypted)\nClient KB Documents", Inches(5.0)),
    ("Bedrock (Claude)\n+ OpenSearch Vectors", Inches(8.5)),
]

for label, x in data_boxes:
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.9), Inches(3.2), Inches(1.1))
    s.fill.solid()
    s.fill.fore_color.rgb = CALLOUT_BG
    s.line.color.rgb = ACCENT_ORANGE
    s.line.width = Pt(1)
    tx = sl.shapes.add_textbox(x + Inches(0.2), Inches(3.95), Inches(2.8), Inches(1.0))
    tf = tx.text_frame
    tf.word_wrap = True
    for j, line in enumerate(label.split("\n")):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(14)
        r.font.color.rgb = WHITE if j == 0 else MED_TEXT
        r.font.bold = (j == 0)
        r.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

# Security footer
txt(sl, "CloudWatch  |  GuardDuty  |  CloudTrail  |  KMS  |  Secrets Manager  |  Stripe",
    Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.4),
    size=13, color=MED_TEXT, align=PP_ALIGN.CENTER)

callout(sl, "All traffic in private subnets. All data encrypted at rest + in transit. Zero public DB access.",
        Inches(0.5), Inches(5.9), Inches(11), color=TEAL)


# ============================================================
# SLIDE 10: Auth + Multi-Tenancy
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
bar(sl)
title_text(sl, "Auth & Multi-Tenancy")
page_num(sl, 10, TOTAL)

# Left: Auth
txt(sl, "Authentication & RBAC", Inches(0.8), Inches(1.2), Inches(5), Inches(0.4),
    size=20, color=TEAL, bold=True)

auth = [
    "Cognito: MFA (required for PHI), SSO (SAML/OIDC), user pools",
    "JWT: 15-min access tokens, 7-day rotating refresh",
    "5 roles: Client Admin, Client User, VA, MEDVA Corporate, MEDVA Admin",
    "Defense in depth: App-level RBAC + DB-level RLS (two independent layers)",
]
bullets(sl, auth, Inches(0.8), Inches(1.7), Inches(5.5), Inches(2.5), size=14)

# Right: Multi-tenancy
txt(sl, "Multi-Tenancy Strategy", Inches(7.0), Inches(1.2), Inches(5), Inches(0.4),
    size=20, color=TEAL, bold=True)

mt = [
    "Shared DB, logical isolation via organization_id + PostgreSQL RLS",
    "RLS enforced at DB level -- even buggy app code can't leak tenant data",
    "Why not separate DBs? 1,000+ clients = unmanageable migrations + no cross-tenant analytics",
]
bullets(sl, mt, Inches(7.0), Inches(1.7), Inches(5.5), Inches(2.0), size=14)

callout(sl, "Even if application logic has a bug, the database won't return another tenant's data.",
        Inches(0.5), Inches(4.2), Inches(11), color=ACCENT_GREEN)


# ============================================================
# SLIDE 11: AI KB Isolation
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
bar(sl)
title_text(sl, "AI KB Isolation")
subtitle(sl, "How We Guarantee No Cross-Tenant Leakage (Highest HIPAA Risk Area)")
page_num(sl, 11, TOTAL)

steps = [
    ("1", "Client uploads document \u2192 S3 under tenant_id/ prefix, encrypted via KMS"),
    ("2", "Async pipeline (SQS + Lambda): parse, chunk, embed"),
    ("3", "Embeddings stored in OpenSearch with org_id metadata tag"),
    ("4", "VA asks question \u2192 OpenSearch query with mandatory org_id filter (infrastructure-level)"),
    ("5", "Retrieved chunks + question sent to Claude via Bedrock (stateless -- no data persists)"),
    ("6", "Response returned with source citations from that client's KB only"),
    ("7", "Nightly automated cross-tenant tests -- must return zero results or pipeline stops"),
]

for i, (num, text) in enumerate(steps):
    y = Inches(1.7) + Inches(i * 0.65)
    # Number circle
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.45), Inches(0.45))
    s.fill.solid()
    s.fill.fore_color.rgb = TEAL
    s.line.fill.background()
    txt(sl, num, Inches(0.8), y + Inches(0.02), Inches(0.45), Inches(0.4),
        size=16, color=DARK_BG, bold=True, align=PP_ALIGN.CENTER)
    txt(sl, text, Inches(1.5), y + Inches(0.03), Inches(10.5), Inches(0.4),
        size=15, color=LIGHT_TEXT)

callout(sl, "HIPAA-safe by design. Bedrock is stateless. Isolation is enforced at infrastructure level, not application level.",
        Inches(0.5), Inches(6.4), Inches(11), color=ACCENT_GREEN)


# ============================================================
# SLIDE 12: Budget
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
bar(sl)
title_text(sl, "Budget & Vendor Management")
page_num(sl, 12, TOTAL)

# Cloud costs
txt(sl, "Monthly Cloud Infrastructure (10-20% adoption)", Inches(0.8), Inches(1.15), Inches(8), Inches(0.35),
    size=18, color=TEAL, bold=True)

rows = [
    ["Phase", "Est. Monthly Cost", "Key Drivers"],
    ["MVP", "$800 - $1,200", "Aurora Serverless (scales to zero), Fargate, S3, CloudFront"],
    ["V1", "$2,200 - $2,900", "+ Bedrock inference ($500-$1,200 variable), OpenSearch, Lambda"],
    ["Full Adoption", "$8K - $12K", "Linear scaling with Fargate auto-scale + Bedrock usage"],
]

table(sl, rows, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.8),
      col_widths=[Inches(1.8), Inches(2.5), Inches(8.0)], size=14)

# Team + Software side by side
txt(sl, "Annual Team Cost (4 hires): ~$330K", Inches(0.8), Inches(3.7), Inches(5), Inches(0.35),
    size=16, color=WHITE, bold=True)
txt(sl, "Software Licensing: ~$400/mo", Inches(7.0), Inches(3.7), Inches(5), Inches(0.35),
    size=16, color=WHITE, bold=True)

txt(sl, "US DevOps ~$195K | Lead Backend ~$52K | Full-Stack ~$46K | UI/UX ~$37K",
    Inches(0.8), Inches(4.15), Inches(6), Inches(0.4), size=13, color=MED_TEXT)
txt(sl, "GitHub, Linear, Figma, Notion, Slack, Claude Code, 1Password, Sentry, PagerDuty. Stripe fees pass-through.",
    Inches(7.0), Inches(4.15), Inches(5.5), Inches(0.6), size=13, color=MED_TEXT)

# Optimization
txt(sl, "Cost Optimization", Inches(0.8), Inches(5.0), Inches(4), Inches(0.35),
    size=18, color=TEAL, bold=True)

opt = [
    "Aurora Serverless: auto-scales to near-zero off-hours",
    "Savings Plans: 1-year compute after MVP (~35% savings)",
    "Fargate Spot: batch processing at up to 70% discount",
    "Bedrock: per-client quotas, prompt caching, cost-efficient tier first",
    "Monthly FinOps reviews with budget alerts at 80% and 100%",
]
bullets(sl, opt, Inches(0.8), Inches(5.4), Inches(11), Inches(1.8), size=14, color=MED_TEXT)


# ============================================================
# SLIDE 13: Build vs Buy
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)
bar(sl)
title_text(sl, "Build vs. Buy")
page_num(sl, 13, TOTAL)

rows = [
    ["Capability", "Decision", "Why"],
    ["Auth", "BUY (Cognito)", "Mature + BAA-covered. Don't reinvent with 5 people."],
    ["Messaging", "BUILD", "Core differentiator + E2E encryption. SendBird scales poorly at 5K+ VAs."],
    ["Medva Frontline AI", "BUILD RAG + Claude via Bedrock", "We own the pipeline + isolation. Strongest model for grounding responses in retrieved docs."],
    ["Medva Academy", "BUILD (V1)", "Mock-up shows deep portal integration. Thinkific can't deliver this."],
    ["AI Extensions", "BUILD framework", "Pluggable architecture. Third-party vendors provide AI via API."],
    ["IT Ticketing", "INTEGRATE", "Not a differentiator. Zendesk/Freshdesk via API."],
    ["Billing", "BUY (Stripe)", "PCI compliance out of the box. Never build payment processing."],
    ["CRM", "INTEGRATE (HubSpot)", "Already in use. Bi-directional sync, don't replace."],
]

table(sl, rows, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.5),
      col_widths=[Inches(2.2), Inches(3.0), Inches(7.1)], size=14)

callout(sl, "Build if it's a core differentiator or touches PHI. Buy if it's commodity. Integrate if MEDVA already uses it.",
        Inches(0.5), Inches(6.2), Inches(11), color=TEAL)


# ============================================================
# SLIDE 14: Closing
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
bg(sl)

# Assumptions
txt(sl, "Key Assumptions", Inches(0.8), Inches(0.6), Inches(8), Inches(0.4),
    size=20, color=TEAL, bold=True)

assumptions = [
    "Building from scratch (no legacy dependencies to inherit)",
    "AWS account with active BAA available or will be provisioned",
    "API/data access to current systems available for migration",
    "Offshore hiring leverages MEDVA's existing Philippines HR infrastructure",
    "10-20% adoption (100-200 clients) for cost modeling",
]
bullets(sl, assumptions, Inches(0.8), Inches(1.1), Inches(11), Inches(2.5), size=15, color=MED_TEXT)

s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.04))
s.fill.solid()
s.fill.fore_color.rgb = TEAL
s.line.fill.background()

txt(sl, "Thank You", Inches(0.8), Inches(3.8), Inches(10), Inches(0.8),
    size=40, color=WHITE, bold=True)

txt(sl, "Happy to go deeper on any section in the walkthrough.",
    Inches(0.8), Inches(4.7), Inches(10), Inches(0.5), size=18, color=MED_TEXT)

txt(sl, "Jimmy Rhoades", Inches(0.8), Inches(5.6), Inches(6), Inches(0.4),
    size=20, color=WHITE, bold=True)
txt(sl, "jimmyrhoades1@gmail.com  |  770.331.7637", Inches(0.8), Inches(6.0), Inches(6), Inches(0.4),
    size=14, color=MED_TEXT)

page_num(sl, 14, TOTAL)


# ============================================================
# Save
# ============================================================
out_dir = os.path.dirname(os.path.abspath(__file__))
out_path = os.path.join(out_dir, "MEDVA_Pulse_Portal_Proposal.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
